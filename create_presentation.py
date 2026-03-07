#!/usr/bin/env python3
"""
Script de génération du PowerPoint de présentation CiSaMe.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Couleurs du thème
BLEU_UNISTRA = RGBColor(0x00, 0x3D, 0x6B)
BLEU_CLAIR = RGBColor(0x4A, 0x90, 0xD9)
GRIS_FONCE = RGBColor(0x33, 0x33, 0x33)
GRIS_CLAIR = RGBColor(0x66, 0x66, 0x66)
BLANC = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xE8, 0x7C, 0x2A)
VERT = RGBColor(0x2E, 0x7D, 0x32)
ROUGE_DOUX = RGBColor(0xC6, 0x28, 0x28)
FOND_CLAIR = RGBColor(0xF5, 0xF7, 0xFA)
FOND_SECTION = RGBColor(0x00, 0x2B, 0x4E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=GRIS_FONCE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(slide, items, left, top, width, height,
                              font_size=16, color=GRIS_FONCE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_numbered_items(slide, items, left, top, width, height,
                        font_size=15, color=GRIS_FONCE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = item
    return txBox


def slide_header(slide, title, subtitle=None):
    """Add a consistent header bar to content slides."""
    add_shape_bg(slide, 0, 0, W, Inches(1.2), BLEU_UNISTRA)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.6),
                 title, font_size=28, color=BLANC, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.4),
                     subtitle, font_size=16, color=RGBColor(0xBB, 0xCC, 0xDD))
    # Thin accent line
    add_shape_bg(slide, 0, Inches(1.2), W, Inches(0.04), ORANGE)


def section_slide(slide, number, title, subtitle=""):
    """Create a dark section divider slide."""
    add_background(slide, FOND_SECTION)
    add_shape_bg(slide, 0, Inches(2.8), W, Inches(2.2), RGBColor(0x00, 0x22, 0x3D))
    add_text_box(slide, Inches(1), Inches(2.0), Inches(2), Inches(0.8),
                 number, font_size=48, color=ORANGE, bold=True)
    add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
                 title, font_size=36, color=BLANC, bold=True)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.9), Inches(11), Inches(0.6),
                     subtitle, font_size=18, color=RGBColor(0x99, 0xAA, 0xBB))


def add_card(slide, left, top, width, height, title, content_lines,
             accent_color=BLEU_CLAIR):
    """Add a card-style box with title and bullet content."""
    # Card background
    card = add_shape_bg(slide, left, top, width, height, BLANC)
    # Accent bar on top
    add_shape_bg(slide, left, top, width, Inches(0.06), accent_color)
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15),
                 width - Inches(0.4), Inches(0.4),
                 title, font_size=16, color=BLEU_UNISTRA, bold=True)
    # Content
    add_bullet_slide_content(slide, content_lines,
                              left + Inches(0.2), top + Inches(0.55),
                              width - Inches(0.4), height - Inches(0.7),
                              font_size=13, color=GRIS_FONCE, spacing=Pt(4))


# =============================================================================
# SLIDE 1 : Page de titre
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide, BLEU_UNISTRA)
add_shape_bg(slide, 0, Inches(1.5), W, Inches(4.5), RGBColor(0x00, 0x34, 0x5C))

add_text_box(slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.5),
             "Universite de Strasbourg", font_size=18, color=RGBColor(0x99, 0xBB, 0xDD))

add_text_box(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
             "CiSaMe", font_size=60, color=BLANC, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(0.8),
             "Circulation des Savoirs Medievaux", font_size=30, color=ORANGE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.8),
             "Projet de recherche en Humanites Numeriques\n"
             "et son integration dans l'environnement des Humanites Numeriques Alsaciennes",
             font_size=18, color=RGBColor(0xCC, 0xDD, 0xEE),
             alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(5.5), Inches(5.2), Inches(2.3), Inches(0.04), ORANGE)

add_text_box(slide, Inches(1), Inches(5.6), Inches(11.3), Inches(0.4),
             "Faculte de Droit | MISHA | Strasbourg",
             font_size=16, color=RGBColor(0x88, 0xAA, 0xCC), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.4),
             "2026",
             font_size=14, color=RGBColor(0x77, 0x99, 0xBB), alignment=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 2 : Sommaire
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, FOND_CLAIR)
slide_header(slide, "Sommaire", "Plan de la presentation (~30 min)")

items = [
    ("01", "Contexte et objectifs du projet CiSaMe", "5 min"),
    ("02", "Architecture technique : le pipeline en 9 modules", "8 min"),
    ("03", "Enrichissement linguistique et validation", "5 min"),
    ("04", "Diffusion et outils de requetage", "5 min"),
    ("05", "Integration dans les Humanites Numeriques Alsaciennes", "4 min"),
    ("06", "Bilan, perspectives et discussion", "3 min"),
]

for i, (num, title, dur) in enumerate(items):
    y = Inches(1.6) + Inches(i * 0.85)
    add_shape_bg(slide, Inches(1), y, Inches(11.3), Inches(0.7), BLANC)
    add_text_box(slide, Inches(1.3), y + Inches(0.12), Inches(0.6), Inches(0.45),
                 num, font_size=22, color=ORANGE, bold=True)
    add_text_box(slide, Inches(2.2), y + Inches(0.12), Inches(7.5), Inches(0.45),
                 title, font_size=18, color=GRIS_FONCE)
    add_text_box(slide, Inches(10.5), y + Inches(0.12), Inches(1.5), Inches(0.45),
                 dur, font_size=14, color=GRIS_CLAIR, alignment=PP_ALIGN.RIGHT)

# =============================================================================
# SECTION 1 : Contexte et objectifs
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(slide, "01", "Contexte et objectifs du projet CiSaMe",
              "De la source medievale au corpus numerique")

# --- Slide 4 : Contexte scientifique ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, FOND_CLAIR)
slide_header(slide, "Contexte scientifique")

add_bullet_slide_content(slide, [
    "Le droit medieval : un patrimoine textuel considerable encore largement inexploite par les outils numeriques",
    "Des centaines de manuscrits juridiques (droit canonique, droit romain) disperses dans les bibliotheques europeennes",
    "Une tradition de recherche riche a l'Universite de Strasbourg (Faculte de Droit, MISHA)",
    "Le besoin d'une plateforme integree pour la numerisation, la transcription et l'analyse computationnelle",
    "CiSaMe : un projet ne a l'intersection de l'histoire du droit et des humanites numeriques",
], Inches(1), Inches(1.6), Inches(11.3), Inches(5),
   font_size=19, color=GRIS_FONCE, spacing=Pt(16))

# --- Slide 5 : Objectifs du projet ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, FOND_CLAIR)
slide_header(slide, "Objectifs du projet CiSaMe")

objectives = [
    ("Numeriser", "Acquerir les images de manuscrits et editions via IIIF, PDF, sources numeriques"),
    ("Transcrire", "Appliquer la reconnaissance automatique d'ecriture manuscrite (HTR) via eScriptorium/Kraken"),
    ("Enrichir", "Annoter linguistiquement les textes latins (lemmatisation, etiquetage morphosyntaxique)"),
    ("Valider", "Verifier la qualite des transcriptions avec un systeme de double dictionnaire"),
    ("Diffuser", "Mettre a disposition les corpus via Nakala, NoSketch-Engine et une plateforme web"),
]

for i, (title, desc) in enumerate(objectives):
    y = Inches(1.6) + Inches(i * 1.05)
    add_shape_bg(slide, Inches(1), y, Inches(0.08), Inches(0.8), ORANGE)
    add_text_box(slide, Inches(1.3), y, Inches(2.5), Inches(0.4),
                 title, font_size=20, color=BLEU_UNISTRA, bold=True)
    add_text_box(slide, Inches(1.3), y + Inches(0.4), Inches(10.5), Inches(0.4),
                 desc, font_size=16, color=GRIS_FONCE)

# --- Slide 6 : Le corpus en chiffres ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, FOND_CLAIR)
slide_header(slide, "Le corpus en chiffres")

stats = [
    ("317", "Manuscrits\nmedievaux repertories"),
    ("129", "Editions\ndocumentees (Heurist)"),
    ("5 768", "Enregistrements\ndans la base Heurist"),
    ("4 149", "Fichiers du\nDecret de Gratien"),
]

for i, (num, label) in enumerate(stats):
    x = Inches(0.8) + Inches(i * 3.15)
    card = add_shape_bg(slide, x, Inches(1.8), Inches(2.8), Inches(2.2), BLANC)
    add_text_box(slide, x, Inches(2.0), Inches(2.8), Inches(0.7),
                 num, font_size=44, color=ORANGE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.8), Inches(2.4), Inches(1.0),
                 label, font_size=15, color=GRIS_FONCE,
                 alignment=PP_ALIGN.CENTER)

# Types de textes
types_data = [
    ("Droit canonique", "~38 editions"),
    ("Theologie", "~61 editions"),
    ("Droit romain", "~17 editions"),
    ("Autres", "~13 editions"),
]

add_text_box(slide, Inches(1), Inches(4.5), Inches(5), Inches(0.4),
             "Repartition par type de texte", font_size=18, color=BLEU_UNISTRA, bold=True)

for i, (typ, count) in enumerate(types_data):
    y = Inches(5.1) + Inches(i * 0.5)
    add_text_box(slide, Inches(1.3), y, Inches(3), Inches(0.4),
                 f"{typ}  :  {count}", font_size=15, color=GRIS_FONCE)

# Depots
add_text_box(slide, Inches(6.5), Inches(4.5), Inches(6), Inches(0.4),
             "Sources principales", font_size=18, color=BLEU_UNISTRA, bold=True)
depots = ["Bibliotheque Vaticane, Archives Secretes du Vatican",
          "British Library, Bodleian Library (Oxford)",
          "Bibliotheque Nationale de France (BnF)",
          "Admont, Bamberg, manuscripta.at, BVMM (IRHT-CNRS)"]
add_bullet_slide_content(slide, depots,
                          Inches(6.5), Inches(5.0), Inches(6), Inches(2.5),
                          font_size=14, color=GRIS_FONCE, spacing=Pt(6))

# =============================================================================
# SECTION 2 : Architecture technique
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(slide, "02", "Architecture technique",
              "Un pipeline modulaire en 9 etapes")

# --- Slide 8 : Vue d'ensemble des 9 modules ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, FOND_CLAIR)
slide_header(slide, "Les 9 modules du pipeline CiSaMe",
             "De l'acquisition des images a la diffusion des corpus annotes")

modules = [
    ("M1", "Telechargement\nimages", BLEU_CLAIR),
    ("M2", "Acquisition\n(IIIF, PDF...)", BLEU_CLAIR),
    ("M3", "Recuperation\neditions", BLEU_CLAIR),
    ("M4", "eScriptorium\nHTR/OCR", ORANGE),
    ("M5", "Nettoyage\nXML", ORANGE),
    ("M6", "Enrichissement\nlinguistique", VERT),
    ("M7", "NoSketch\nEngine", VERT),
    ("M8", "Diffusion\nNakala/Seafile", RGBColor(0x7B, 0x1F, 0xA2)),
    ("M9", "Visualisation\nRequetes", RGBColor(0x7B, 0x1F, 0xA2)),
]

for i, (code, label, color) in enumerate(modules):
    col = i % 5
    row = i // 5
    x = Inches(0.6) + Inches(col * 2.55)
    y = Inches(1.8) + Inches(row * 2.6)

    box = add_shape_bg(slide, x, y, Inches(2.2), Inches(2.0), BLANC)
    add_shape_bg(slide, x, y, Inches(2.2), Inches(0.06), color)
    add_text_box(slide, x, y + Inches(0.2), Inches(2.2), Inches(0.4),
                 code, font_size=20, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.7), Inches(2.0), Inches(1.1),
                 label, font_size=14, color=GRIS_FONCE,
                 alignment=PP_ALIGN.CENTER)

    # Arrow between boxes (except last in row)
    if col < 4 and i < 8:
        add_text_box(slide, x + Inches(2.2), y + Inches(0.7), Inches(0.35), Inches(0.4),
                     ">", font_size=24, color=GRIS_CLAIR, bold=True,
                     alignment=PP_ALIGN.CENTER)

# Transversal bar
add_shape_bg(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.6), BLEU_UNISTRA)
add_text_box(slide, Inches(1), Inches(6.55), Inches(11), Inches(0.4),
             "Transversal  :  Base de metadonnees Heurist (hdb_cisame_misha)  |  "
             "Pipeline parallele Decret de Gratien",
             font_size=14, color=BLANC)

# --- Slide 9 : Acquisition et transcription ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, FOND_CLAIR)
slide_header(slide, "Modules 1-4 : Acquisition et Transcription")

add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.5),
         "Acquisition des sources (M1-M3)", [
             "Telechargement via protocole IIIF (manifestes JSON)",
             "Extraction d'images depuis PDF, hexadecimal, tuiles",
             "Classification par statut de droits : libre / restreint / prive",
             "Inventaire de 317 manuscrits dans liste_manuscrits.csv",
         ], BLEU_CLAIR)

add_card(slide, Inches(7), Inches(1.6), Inches(5.8), Inches(2.5),
         "Transcription HTR/OCR (M4)", [
             "Plateforme eScriptorium pour segmentation + transcription",
             "Modeles Kraken entraines et fine-tunes (HPC Unistra)",
             "Modeles distincts : manuscrits vs editions imprimees",
             "Export au format XML PAGE (standard HTR)",
         ], ORANGE)

add_card(slide, Inches(0.5), Inches(4.5), Inches(5.8), Inches(2.5),
         "Nettoyage XML (M5)", [
             "Nettoyage via Oxygen XML Editor (XPath + Regex)",
             "Suppression des artefacts de segmentation",
             "Verification de la structure XML PAGE",
             "Preparation pour l'enrichissement linguistique",
         ], ORANGE)

add_card(slide, Inches(7), Inches(4.5), Inches(5.8), Inches(2.5),
         "Technologies cles", [
             "Kraken : moteur ML pour HTR/OCR",
             "eScriptorium : interface web de transcription",
             "XML PAGE : standard d'echange pour la transcription",
             "Fine-tuning sur le HPC de l'Universite de Strasbourg",
         ], GRIS_CLAIR)

# --- Slide 10 : Modeles Kraken ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, FOND_CLAIR)
slide_header(slide, "Modeles Kraken : segmentation et transcription")

add_card(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(2.3),
         "Modeles de segmentation", [
             "SEGMENTATION_EDITION : detection de lignes pour editions imprimees",
             "SEGMENTATION_MANUSCRIT : detection pour ecritures manuscrites",
             "Identification automatique des zones (MainZone, MarginZone...)",
         ], BLEU_CLAIR)

add_card(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(2.3),
         "Modeles de transcription", [
             "TRANSCRIPTION_EDITION : reconnaissance de caracteres imprimes",
             "TRANSCRIPTION_MANUSCRIT : reconnaissance d'ecritures medievales",
             "Fine-tuning iteratif pour ameliorer la precision",
         ], ORANGE)

add_card(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.7),
         "Processus de fine-tuning (HPC Unistra)", [
             "1. Transcription manuelle d'un echantillon de pages (ground truth)",
             "2. Entrainement sur le cluster HPC de l'Universite de Strasbourg",
             "3. Evaluation du Character Error Rate (CER) et du Word Error Rate (WER)",
             "4. Iteration : correction des erreurs, re-entrainement, amelioration progressive",
             "5. Partage des modeles via HTR-United (registre communautaire)",
         ], VERT)

# =============================================================================
# SECTION 3 : Enrichissement linguistique et validation
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(slide, "03", "Enrichissement linguistique et validation",
              "PAGEtopage et latin_analyzer")

# --- Slide 12 : PAGEtopage pipeline ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, FOND_CLAIR)
slide_header(slide, "Module 6 : PAGEtopage - Pipeline d'enrichissement",
             "Du XML PAGE au corpus annote")

steps = [
    ("Step 1\nExtract", "Parsing XML PAGE\nExtraction du texte\nFusion des mots coupes", BLEU_CLAIR),
    ("Step 2\nEnrich", "Tokenisation\nTreeTagger (lemmes)\nEtiquetage POS", ORANGE),
    ("Step 3\nExport", "4 formats de sortie :\nscholarly, clean,\ndiplomatic, annotated", VERT),
    ("Step 4\nRe-enrich", "Corrections manuelles\nRe-generation auto\ndes annotations", RGBColor(0x7B, 0x1F, 0xA2)),
    ("Step 5\nDOCX>Vert.", "Conversion DOCX\nvers format vertical\n(latin_analyzer)", ROUGE_DOUX),
]

for i, (title, desc, color) in enumerate(steps):
    x = Inches(0.4) + Inches(i * 2.6)
    box = add_shape_bg(slide, x, Inches(1.8), Inches(2.3), Inches(3.2), BLANC)
    add_shape_bg(slide, x, Inches(1.8), Inches(2.3), Inches(0.06), color)
    add_text_box(slide, x, Inches(2.0), Inches(2.3), Inches(0.8),
                 title, font_size=16, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(2.9), Inches(2.0), Inches(1.8),
                 desc, font_size=13, color=GRIS_FONCE,
                 alignment=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(slide, x + Inches(2.3), Inches(2.8), Inches(0.3), Inches(0.5),
                     ">", font_size=28, color=GRIS_CLAIR, bold=True)

# Example annotation
add_shape_bg(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.8), BLANC)
add_text_box(slide, Inches(0.7), Inches(5.4), Inches(5), Inches(0.4),
             "Exemple d'annotation TreeTagger (format vertical) :", font_size=14,
             color=BLEU_UNISTRA, bold=True)
example_text = (
    "Mot          POS        Lemme\n"
    "Dominus      NOM        dominus\n"
    "enim         ADV        enim\n"
    "dicit        VER        dico"
)
add_text_box(slide, Inches(0.7), Inches(5.85), Inches(5), Inches(1.2),
             example_text, font_size=13, color=GRIS_FONCE, font_name="Consolas")

add_text_box(slide, Inches(7), Inches(5.4), Inches(5), Inches(0.4),
             "Technologies :", font_size=14, color=BLEU_UNISTRA, bold=True)
add_bullet_slide_content(slide, [
    "TreeTagger : lemmatisation et POS-tagging",
    "PyCollatinus : dictionnaire latin classique (~500k formes)",
    "Du Cange : dictionnaire latin medieval (~100k entrees)",
    "Gestion des variantes orthographiques (u/v, i/j, ae/e)",
], Inches(7), Inches(5.8), Inches(5.5), Inches(1.3),
   font_size=13, color=GRIS_FONCE, spacing=Pt(4))

# --- Slide 13 : Formats d'export ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, FOND_CLAIR)
slide_header(slide, "Les 4 formats d'export")

formats = [
    ("Scholarly", "Format academique complet avec en-tetes de metadonnees et texte continu. "
     "Destine a la publication et a la citation scientifique.", BLEU_CLAIR),
    ("Clean", "Texte brut, maximalement lisible. Sans annotations ni metadonnees. "
     "Pour la lecture et la reference rapide.", VERT),
    ("Diplomatic", "Texte avec annotations POS en ligne. Conserve les informations "
     "linguistiques dans un format lisible.", ORANGE),
    ("Annotated", "Format tabulaire pour l'analyse computationnelle. "
     "Chaque mot sur une ligne avec lemme et etiquette POS.", RGBColor(0x7B, 0x1F, 0xA2)),
]

for i, (name, desc, color) in enumerate(formats):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(col * 6.4)
    y = Inches(1.6) + Inches(row * 2.7)
    add_card(slide, x, y, Inches(6), Inches(2.3),
             f"Format {name}", [desc], color)

# --- Slide 14 : latin_analyzer ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, FOND_CLAIR)
slide_header(slide, "Module 5B : latin_analyzer - Validation des transcriptions",
             "Systeme de scoring automatique v2.4.0")

# Score ranges
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.4),
             "Systeme de scoring (0-100 points)", font_size=20, color=BLEU_UNISTRA, bold=True)

scores = [
    ("75-100", "Mot reconnu, probablement correct", RGBColor(0x1B, 0x5E, 0x20), "Noir"),
    ("40-74", "Douteux, verification manuelle recommandee", ORANGE, "Orange"),
    ("0-39", "Erreur OCR probable", ROUGE_DOUX, "Rouge"),
]

for i, (rang, desc, color, label) in enumerate(scores):
    y = Inches(2.2) + Inches(i * 0.7)
    add_shape_bg(slide, Inches(0.8), y, Inches(0.15), Inches(0.5), color)
    add_text_box(slide, Inches(1.2), y, Inches(1.5), Inches(0.5),
                 rang, font_size=18, color=color, bold=True)
    add_text_box(slide, Inches(2.8), y + Inches(0.05), Inches(4), Inches(0.4),
                 f"({label}) {desc}", font_size=15, color=GRIS_FONCE)

# Sources de validation
add_text_box(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.4),
             "Sources de validation", font_size=20, color=BLEU_UNISTRA, bold=True)

sources = [
    "1. PyCollatinus (latin classique, ~500 000 formes)",
    "2. Du Cange (latin medieval, ~100 000 entrees)",
    "3. Suffixes latins productifs (-arius, -atio, -torium)",
    "4. Detection de contexte ecclesiastique",
    "5. Normalisation des variantes (u/v, i/j, ae/e)",
]
add_numbered_items(slide, sources,
                    Inches(7), Inches(2.2), Inches(5.5), Inches(3),
                    font_size=15, color=GRIS_FONCE)

add_card(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2),
         "Sortie : document DOCX avec code couleur", [
             "Chaque mot du texte est colore selon son score de fiabilite",
             "Rapport statistique genere automatiquement (% de mots reconnus, douteux, errones)",
             "Permet une correction manuelle ciblee : le chercheur se concentre sur les zones rouges/orange",
             "Integration avec le pipeline PAGEtopage pour re-enrichissement apres corrections (Step 4)",
         ], VERT)

# =============================================================================
# SECTION 4 : Diffusion et outils de requetage
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(slide, "04", "Diffusion et outils de requetage",
              "NoSketch-Engine, Nakala, Canon-Law-Toolkit")

# --- Slide 16 : Module 7 - NoSketch-Engine ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, FOND_CLAIR)
slide_header(slide, "Module 7 : NoSketch-Engine - Concordancier",
             "Requetes linguistiques sur le corpus CiSaMe")

add_card(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(2.5),
         "Fonctionnalites", [
             "Corpus Query Language (CQL) pour des requetes avancees",
             "Recherche par mot, lemme, categorie grammaticale",
             "Concordances KWIC (Key Word In Context)",
             "Frequences, collocations, n-grammes",
             "Comparaison entre corpus et sous-corpus",
         ], BLEU_CLAIR)

add_card(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(2.5),
         "Infrastructure", [
             "Instance hebergee a l'Universite de Strasbourg",
             "Serveur : fip-185-155-93-80.iaas.unistra.fr",
             "Format d'entree : fichiers verticaux (mot/POS/lemme)",
             "Fusion automatique des fichiers (fusion_vertical.py)",
         ], ORANGE)

add_card(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
         "Module 7 : Processus de mise en corpus", [
             "1. Generation des fichiers verticaux par PAGEtopage (Module 6, format annote)",
             "2. Fusion des fichiers verticaux d'une edition (fusion_vertical.py)",
             "3. Ajout des metadonnees Heurist dans les en-tetes (Edi-XX, auteur, date, type...)",
             "4. Chargement dans NoSketch-Engine et compilation du corpus",
             "5. Le corpus est interrogeable en CQL par les chercheurs",
         ], VERT)

# --- Slide 17 : Module 8 - Diffusion ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, FOND_CLAIR)
slide_header(slide, "Module 8 : Diffusion des donnees",
             "Nakala (public) et Seafile (acces restreint)")

add_card(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(3.5),
         "Nakala (Huma-Num)", [
             "Depot scientifique avec identifiants perennes (DOI)",
             "Upload automatise via l'API Heimdall",
             "Editions libres de droits uniquement",
             "Scripts : validate_export.py > prepare_nakala_export.py > upload_nakala.py",
             "Enrichissement des verticaux avec URLs Nakala",
             "Conversion des fiches de metadonnees en PDF (LibreOffice)",
         ], BLEU_CLAIR)

add_card(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(3.5),
         "Seafile (Cloud prive)", [
             "Stockage securise pour les materiaux sous droits",
             "Acces restreint aux membres du projet",
             "Editions sous copyright ou acces limite",
             "Complementaire a Nakala pour la couverture complete",
         ], ORANGE)

add_card(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5),
         "Gestion des droits : un triple classement", [
             "Libre : diffusion ouverte sur Nakala avec DOI  |  "
             "Restreint : acces limite via Seafile  |  "
             "Prive : acces exclusif equipe projet",
         ], VERT)

# --- Slide 18 : Module 9 - Canon-Law-Toolkit ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, FOND_CLAIR)
slide_header(slide, "Module 9 : Canon-Law-Toolkit",
             "Plateforme web d'analyse - cisame.vercel.app")

add_card(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(2.3),
         "Generateur de requetes CQL", [
             "Recherches de proximite (1-20 tokens)",
             "96 variantes orthographiques automatiques",
             "Recherches semantiques par lemme",
             "Combinaisons avancees proximite + variations",
         ], BLEU_CLAIR)

add_card(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(2.3),
         "Analyseur de concordances", [
             "9 vues specialisees pour l'analyse de corpus",
             "Enrichissement automatique avec metadonnees Edi-XX",
             "Visualisations interactives (Recharts, D3.js)",
             "Export multi-format : CSV, JSON, PNG",
         ], ORANGE)

add_card(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.8),
         "Technologies et deploiement", [
             "Frontend : React 18.2 + Vite 5.0 (build rapide)",
             "Visualisation : Recharts + D3.js (graphiques, timelines, nuages de mots)",
             "Interface bilingue FR/EN",
             "Deploiement CI/CD sur Vercel (mise a jour automatique)",
             "Version 1.5.0 (novembre 2025) - tests : 93/93 composants UI (100%)",
         ], VERT)

# =============================================================================
# SECTION 5 : Integration HN Alsaciennes
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(slide, "05", "Integration dans les Humanites Numeriques Alsaciennes",
              "Un projet ancre dans l'ecosysteme regional et national")

# --- Slide 20 : Ecosysteme HN ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, FOND_CLAIR)
slide_header(slide, "L'ecosysteme des Humanites Numeriques en Alsace")

add_card(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(2.3),
         "Ancrage institutionnel", [
             "Universite de Strasbourg (Unistra)",
             "Faculte de Droit - Histoire du droit medieval",
             "MISHA : Maison Interuniversitaire des Sciences de l'Homme - Alsace",
             "Cluster HPC de l'Universite pour le calcul intensif",
         ], BLEU_CLAIR)

add_card(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(2.3),
         "Infrastructure technique regionale", [
             "Serveur NoSketch-Engine heberge a l'Unistra (IaaS)",
             "HPC pour le fine-tuning des modeles Kraken",
             "Integration avec les services numeriques de l'universite",
             "Hebergement et maintenance par l'equipe technique",
         ], ORANGE)

add_card(slide, Inches(0.5), Inches(4.3), Inches(6), Inches(2.8),
         "Connexions nationales (Huma-Num / TGIR)", [
             "Nakala : depot de donnees scientifiques perennes",
             "Heimdall : API d'upload vers les services Huma-Num",
             "Participation au reseau HTR-United",
             "Utilisation des standards nationaux DH (XML PAGE, CQL...)",
         ], VERT)

add_card(slide, Inches(6.8), Inches(4.3), Inches(6), Inches(2.8),
         "Outils partages de la communaute DH", [
             "eScriptorium : plateforme HTR collaborative",
             "Kraken : moteur open-source de reconnaissance",
             "TreeTagger : etiquetage morphosyntaxique multilingue",
             "NoSketch-Engine : concordancier open-source",
             "Heurist : base de donnees relationnelle pour les SHS",
         ], RGBColor(0x7B, 0x1F, 0xA2))

# --- Slide 21 : Modele de donnees Heurist ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, FOND_CLAIR)
slide_header(slide, "Base de metadonnees Heurist (hdb_cisame_misha)",
             "Modele relationnel a 3 tables principales")

# 3 entity boxes
entities = [
    ("AUTEURS", "Nom, variantes\nmultilingues", "Ex: Bernardus\nPapiensis", BLEU_CLAIR),
    ("OEUVRES", "Titre, date, lieu,\ntype, responsable", "Ex: Summa titulorum\ndecretalium", ORANGE),
    ("EDITIONS", "Editeur, lieu, date,\npages, oeuvre (FK)", "Ex: Laspeyres,\nRegensburg, 1860", VERT),
]

for i, (name, fields, example, color) in enumerate(entities):
    x = Inches(0.8) + Inches(i * 4.2)
    add_shape_bg(slide, x, Inches(1.8), Inches(3.6), Inches(2.8), BLANC)
    add_shape_bg(slide, x, Inches(1.8), Inches(3.6), Inches(0.06), color)
    add_text_box(slide, x, Inches(1.95), Inches(3.6), Inches(0.4),
                 name, font_size=20, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(3.2), Inches(0.8),
                 fields, font_size=14, color=GRIS_FONCE)
    add_text_box(slide, x + Inches(0.2), Inches(3.4), Inches(3.2), Inches(0.8),
                 example, font_size=12, color=GRIS_CLAIR)
    if i < 2:
        add_text_box(slide, x + Inches(3.6), Inches(2.5), Inches(0.6), Inches(0.5),
                     "1:n >", font_size=14, color=ORANGE, bold=True)

add_card(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
         "Statistiques et equipe", [
             "5 768 enregistrements  |  129 editions documentees  |  7 utilisateurs actifs",
             "Responsables : Yann, Raphael Eckert, Guillaume Porte et collaborateurs",
             "6 types d'entites crees (Auteur, Oeuvre, Edition, Commentaire, Allegation, Date)",
             "Integration avec les modules 6, 7 et 8 pour les metadonnees des corpus",
         ], BLEU_UNISTRA)

# =============================================================================
# SECTION 6 : Bilan et perspectives
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(slide, "06", "Bilan, perspectives et discussion",
              "Un projet operationnel ouvert sur l'avenir")

# --- Slide 23 : Bilan ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, FOND_CLAIR)
slide_header(slide, "Bilan du projet")

add_card(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(2.3),
         "Realisations techniques", [
             "9 modules pleinement operationnels",
             "Pipeline complet : de l'image au corpus interrogeable",
             "Plateforme web deployee (cisame.vercel.app)",
             "100% des composants UI testes (93/93)",
         ], VERT)

add_card(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(2.3),
         "Realisations scientifiques", [
             "317 manuscrits repertories et documentes",
             "129 editions integrees dans Heurist",
             "Corpus interrogeable sur NoSketch-Engine",
             "Donnees diffusees sur Nakala (avec DOI)",
         ], BLEU_CLAIR)

add_card(slide, Inches(0.5), Inches(4.3), Inches(6), Inches(2.8),
         "Apports methodologiques", [
             "Chaine de traitement reproductible et documentee",
             "Guide utilisateur de 700+ lignes pour PAGEtopage",
             "9 fichiers README + 4 rapports d'analyse",
             "Modeles Kraken partages via HTR-United",
             "Code source versionne et maintenable",
         ], ORANGE)

add_card(slide, Inches(6.8), Inches(4.3), Inches(6), Inches(2.8),
         "Integration reussie dans l'ecosysteme HN", [
             "Utilisation des standards nationaux (Nakala, Huma-Num)",
             "Infrastructure hebergee a l'Universite de Strasbourg",
             "Outils open-source partages avec la communaute DH",
             "Interface bilingue FR/EN pour rayonnement international",
             "Modele reutilisable pour d'autres corpus medievaux",
         ], RGBColor(0x7B, 0x1F, 0xA2))

# --- Slide 24 : Perspectives ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, FOND_CLAIR)
slide_header(slide, "Perspectives et developpements futurs")

add_bullet_slide_content(slide, [
    "Extension du corpus : integration de nouveaux manuscrits et traditions textuelles",
    "Amelioration des modeles HTR : fine-tuning continu pour reduire le taux d'erreur",
    "Enrichissement de la base Heurist : ajout des entites manquantes (Chapitre, Lien)",
    "Interoperabilite : connexion avec d'autres projets DH regionaux et europeens",
    "Analyse computationnelle avancee : detection automatique d'allegation, analyse de reseau",
    "Formation et documentation : ateliers pour la communaute HN alsacienne",
    "Perennisation : integration dans les infrastructures durables de l'Unistra et Huma-Num",
], Inches(1), Inches(1.6), Inches(11.3), Inches(5),
   font_size=18, color=GRIS_FONCE, spacing=Pt(18))

# --- Slide 25 : Merci / Discussion ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, BLEU_UNISTRA)
add_shape_bg(slide, 0, Inches(2.2), W, Inches(3.5), RGBColor(0x00, 0x34, 0x5C))

add_text_box(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1),
             "Merci de votre attention", font_size=44, color=BLANC, bold=True,
             alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(5.5), Inches(3.6), Inches(2.3), Inches(0.04), ORANGE)

add_text_box(slide, Inches(1), Inches(3.9), Inches(11.3), Inches(0.6),
             "Questions et discussion", font_size=24, color=RGBColor(0xCC, 0xDD, 0xEE),
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
             "CiSaMe - Circulation des Savoirs Medievaux",
             font_size=18, color=RGBColor(0x88, 0xAA, 0xCC),
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(0.5),
             "cisame.vercel.app  |  Universite de Strasbourg  |  Faculte de Droit  |  MISHA",
             font_size=14, color=RGBColor(0x77, 0x99, 0xBB),
             alignment=PP_ALIGN.CENTER)

# =============================================================================
# Sauvegarde
# =============================================================================
output_path = "/home/user/CiSaMe/Presentation_CiSaMe_Humanites_Numeriques_Alsaciennes.pptx"
prs.save(output_path)
print(f"Presentation sauvegardee : {output_path}")
print(f"Nombre de slides : {len(prs.slides)}")
